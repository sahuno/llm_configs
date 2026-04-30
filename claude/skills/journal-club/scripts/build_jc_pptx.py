#!/usr/bin/env python3
"""
build_jc_pptx.py — Convert Stage 5 slide markdown into a PPTX with embedded figures.

Author: Samuel Ahuno (ekwame001@gmail.com)
Date: 2026-04-29

Usage:
    python build_jc_pptx.py <paper_dir>
    python build_jc_pptx.py <paper_dir> --output /custom/path.pptx
    python build_jc_pptx.py <paper_dir> --images /custom/images/dir
    python build_jc_pptx.py <paper_dir> --max-width 1600
    python build_jc_pptx.py <paper_dir> --no-resize

Conventions assumed:
    <paper_dir>/05_slides_draft.md   — Stage 5 slide content (required)
    <paper_dir>/_meta.json           — paper metadata (read for paper_id)
    <paper_dir>/images/              — extracted figure panels (PNG)
    <paper_dir>/pdf/*.pdf            — source PDF (optional; enables figure↔caption pairing)

Output:
    <paper_dir>/<paper_id>.pptx (default; override with --output)

Layout (16:9 widescreen, 13.33"×7.5"):
    Title:          0.5", 0.3" → 12.33"×1.0"
    Text-only body: 0.5", 1.4" → 12.33"×5.6"
    Image-slide bullets: 0.5", 1.4" → 6.2"×5.6"
    Image:          6.9", 1.4" → 5.94"×5.6" (preserves aspect ratio)

Image assignment strategy (in order of preference):
    1. Caption-aware: parse PDF for "Figure N" caption pages; when a slide's
       Visual field references "Figure N", embed the image whose page is
       closest to that caption page (typically caption_page - 1).
    2. Sequential fallback: if no figure mentioned (or no PDF available),
       walk images in page order, one per image slide.
"""

import argparse
import json
import re
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from PIL import Image
except ImportError as e:
    sys.exit(
        f"Missing dependency: {e}\n"
        "Install with:\n"
        "  pip install python-pptx Pillow"
    )


# ---------- markdown parsing ---------------------------------------------------

SLIDE_PATTERN = re.compile(
    r"^## Slide (\d+)\s*[\u2014\-]\s*(.+?)\n(.*?)(?=^## Slide |\Z)",
    re.MULTILINE | re.DOTALL,
)

STOP_MARKERS = (
    "## Pre-talk checklist",
    "## Conversion to PowerPoint",
    "## Source map",
)

IMAGE_KEYWORDS = (
    "figure", "fig ", "fig.", "panel", "cartoon", "schematic",
    "diagram", "map", "cryo-em", "cryo-EM", "side-by-side", "split",
    "overlay", "structure",
)


def parse_field(block: str, label: str) -> str:
    pat = rf"\*\*{label}\*\*:\s*(.*?)(?=\n\*\*[A-Za-z][A-Za-z ]+\*\*:|\n---|\Z)"
    m = re.search(pat, block, re.DOTALL)
    return m.group(1).strip() if m else ""


def parse_bullets(text: str) -> list[str]:
    return [
        line.strip()[1:].strip()
        for line in text.splitlines()
        if line.strip().startswith("-")
    ]


def parse_slides(md_path: Path) -> list[dict]:
    text = md_path.read_text()
    for marker in STOP_MARKERS:
        idx = text.find(marker)
        if idx != -1:
            text = text[:idx]
    slides = []
    for m in SLIDE_PATTERN.finditer(text):
        body = m.group(3)
        notes = parse_field(body, "Speaker notes")
        if notes.startswith('"') and notes.endswith('"'):
            notes = notes[1:-1]
        slides.append({
            "n": int(m.group(1)),
            "title": m.group(2).strip(),
            "bullets": parse_bullets(parse_field(body, "Bullets")),
            "notes": notes,
            "visual": parse_field(body, "Visual"),
        })
    return slides


def needs_image(visual_text: str) -> bool:
    """True if the slide's Visual: field implies a figure should be embedded."""
    if not visual_text:
        return False
    low = visual_text.lower()
    if "text-only" in low:
        return False
    return any(kw in low for kw in IMAGE_KEYWORDS)


# ---------- caption-aware figure mapping ---------------------------------------

CAPTION_RE = re.compile(
    r"(?:^|\n)\s*(?:Figure|Fig\.|FIGURE)\s+(\d+)\b",
    re.IGNORECASE,
)
SLIDE_FIG_REF_RE = re.compile(
    r"\b(?:figure|fig\.?)\s*(\d+)",
    re.IGNORECASE,
)
PAGE_FROM_FILENAME_RE = re.compile(r"page0*(\d+)_img", re.IGNORECASE)


def parse_caption_pages(pdf_path: Path) -> dict[int, int]:
    """
    Scan a PDF for figure caption headings and return the first page where
    each caption appears. Falls back to {} if PyMuPDF isn't installed or the
    PDF can't be read.

    Returns
    -------
    dict mapping {figure_number: 1-indexed_page_number}
    """
    if not pdf_path.exists():
        return {}
    try:
        import fitz  # PyMuPDF
    except ImportError:
        return {}
    fig_pages: dict[int, int] = {}
    try:
        doc = fitz.open(str(pdf_path))
        try:
            for i, page in enumerate(doc, start=1):
                text = page.get_text()
                for m in CAPTION_RE.finditer(text):
                    n = int(m.group(1))
                    # Record only first occurrence of each figure caption
                    fig_pages.setdefault(n, i)
        finally:
            doc.close()
    except Exception:
        return {}
    return fig_pages


def image_page(img_path: Path) -> int | None:
    """Extract page number from a PyMuPDF-style filename like page012_img01.png."""
    m = PAGE_FROM_FILENAME_RE.search(img_path.name)
    return int(m.group(1)) if m else None


def slide_figure_ref(visual_text: str) -> int | None:
    """If the slide's Visual: field references a numbered figure, return the
    figure number; else None."""
    m = SLIDE_FIG_REF_RE.search(visual_text or "")
    return int(m.group(1)) if m else None


def resolve_image_for_figure(
    fig_num: int,
    caption_pages: dict[int, int],
    images: list[Path],
    used: set[Path],
    figure_to_image: dict[int, Path],
) -> Path | None:
    """
    Find the best image for a referenced figure number.

    Priority (in order):
      0. If this figure was already mapped (e.g., earlier slide referenced
         the same Figure N), reuse that image — slides for "Figure 3A" and
         "Figure 3D" both get the Figure 3 page-image, by design.
      1. Look up caption page for fig_num.
      2. Walk candidate pages in *typical layout order*:
         caption_page - 1 (figure renders on prior page; most common)
         → caption_page (figure on the same page as caption)
         → caption_page + 1 (figure overflow to next page).
      3. Skip images already claimed by a different figure.
      4. Pick the image with the largest file size on the chosen page.

    Returns None if no caption is known or no usable image found; caller
    falls back to sequential.
    """
    # Re-use mapping if a previous slide already mapped this figure
    prior = figure_to_image.get(fig_num)
    if prior is not None:
        return prior

    cap_page = caption_pages.get(fig_num)
    if cap_page is None:
        return None

    # Layout-typical page priority
    for offset in (-1, 0, 1):
        page = cap_page + offset
        candidates = [
            p for p in images
            if image_page(p) == page and p not in used
        ]
        if candidates:
            best = max(candidates, key=lambda p: p.stat().st_size)
            used.add(best)
            figure_to_image[fig_num] = best
            return best

    return None


# ---------- pptx building ------------------------------------------------------

def fit_image_box(img_path: Path, max_w_in: float, max_h_in: float
                  ) -> tuple[float, float]:
    """Return (width, height) in inches that fits img into (max_w, max_h)
    while preserving aspect ratio."""
    with Image.open(img_path) as img:
        iw, ih = img.size
    aspect = iw / ih
    if max_w_in / max_h_in > aspect:
        h = max_h_in
        w = h * aspect
    else:
        w = max_w_in
        h = w / aspect
    return w, h


def set_font(run, name: str = "Arial", size: int | None = None,
             bold: bool | None = None) -> None:
    run.font.name = name
    if size is not None:
        run.font.size = Pt(size)
    if bold is not None:
        run.font.bold = bold


def maybe_resize(src_path: Path, max_width: int, tmp_dir: Path) -> Path:
    """
    If src image is wider than max_width pixels, resize (LANCZOS) and save
    to tmp_dir; return the path to use. Otherwise return src_path unchanged.

    max_width=0 disables resizing entirely.
    """
    if max_width <= 0:
        return src_path
    with Image.open(src_path) as img:
        if img.width <= max_width:
            return src_path
        ratio = max_width / img.width
        new_size = (max_width, int(img.height * ratio))
        resized = img.resize(new_size, Image.LANCZOS)
        out = tmp_dir / src_path.name
        # Save as PNG (preserves transparency); optimize for size
        resized.save(out, format="PNG", optimize=True)
    return out


def build_pptx(
    slides: list[dict],
    images_dir: Path,
    out_path: Path,
    pdf_path: Path | None = None,
    max_width: int = 1600,
) -> dict:
    """
    Build the journal-club PPTX.

    Returns a stats dict for the caller to log.
    """
    import tempfile

    prs = Presentation()
    # 16:9 widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # blank — gives full control of shapes

    # Available images sorted by filename (PyMuPDF outputs are page-sorted)
    image_files = sorted(images_dir.glob("*.png")) if images_dir.exists() else []
    images_total = len(image_files)

    # Caption-aware figure mapping (best effort; empty dict = sequential fallback)
    caption_pages: dict[int, int] = {}
    if pdf_path is not None:
        caption_pages = parse_caption_pages(pdf_path)

    # Sequential pool: images we haven't yet placed via figure resolution
    seq_iter = iter(image_files)
    placed_via_caption: set[Path] = set()
    figure_to_image: dict[int, Path] = {}  # cache so Fig 3A and Fig 3D share image
    images_used = 0
    seq_used = 0
    cap_used = 0

    # Temp dir for resized images (kept until the prs is saved)
    tmp_dir = Path(tempfile.mkdtemp(prefix="jcpptx_"))

    # Layout constants (inches)
    TITLE_X, TITLE_Y, TITLE_W, TITLE_H = 0.5, 0.3, 12.33, 1.0
    FULL_BODY_X, FULL_BODY_Y, FULL_BODY_W, FULL_BODY_H = 0.5, 1.4, 12.33, 5.6
    HALF_BODY_X, HALF_BODY_Y, HALF_BODY_W, HALF_BODY_H = 0.5, 1.4, 6.2, 5.6
    IMG_X, IMG_Y, IMG_MAX_W, IMG_MAX_H = 6.9, 1.4, 5.94, 5.6

    def next_sequential_image() -> Path | None:
        """Walk the sequential iterator, skipping any image already placed
        by caption-resolution, until we find an unused one."""
        nonlocal seq_used
        for p in seq_iter:
            if p not in placed_via_caption:
                seq_used += 1
                return p
        return None

    for s in slides:
        slide = prs.slides.add_slide(blank_layout)

        is_title = (s["n"] == 1)
        is_image = (not is_title) and needs_image(s["visual"])

        # ---- Title shape ----
        title_box = slide.shapes.add_textbox(
            Inches(TITLE_X), Inches(TITLE_Y),
            Inches(TITLE_W), Inches(TITLE_H),
        )
        title_tf = title_box.text_frame
        title_tf.word_wrap = True
        if is_title and s["bullets"]:
            paper_title = s["bullets"][0]
        else:
            paper_title = s["title"]
        title_tf.text = paper_title
        for run in title_tf.paragraphs[0].runs:
            set_font(run, "Arial", 32 if is_title else 28, bold=True)

        # ---- Body / image ----
        image_source_note = None  # diagnostic for slide notes

        if is_title:
            # Subtitle: remaining bullets, centered visual region
            subtitle_lines = s["bullets"][1:] if len(s["bullets"]) > 1 else []
            if subtitle_lines:
                sub_box = slide.shapes.add_textbox(
                    Inches(FULL_BODY_X), Inches(FULL_BODY_Y + 1.5),
                    Inches(FULL_BODY_W), Inches(3.5),
                )
                sub_tf = sub_box.text_frame
                sub_tf.word_wrap = True
                for i, line in enumerate(subtitle_lines):
                    para = sub_tf.paragraphs[0] if i == 0 else sub_tf.add_paragraph()
                    para.text = line
                    for run in para.runs:
                        set_font(run, "Arial", 18)
        else:
            # Bullets — full width if text-only, half width if image slide
            if is_image:
                body_x, body_w = HALF_BODY_X, HALF_BODY_W
            else:
                body_x, body_w = FULL_BODY_X, FULL_BODY_W
            body_box = slide.shapes.add_textbox(
                Inches(body_x), Inches(FULL_BODY_Y),
                Inches(body_w), Inches(FULL_BODY_H),
            )
            body_tf = body_box.text_frame
            body_tf.word_wrap = True
            for i, b in enumerate(s["bullets"]):
                para = body_tf.paragraphs[0] if i == 0 else body_tf.add_paragraph()
                para.text = "• " + b
                para.level = 0
                for run in para.runs:
                    set_font(run, "Arial", 18)

            # Embed an image on image slides
            if is_image:
                # Try caption-aware first — check Visual: field, then slide title
                img_path: Path | None = None
                fig_num = slide_figure_ref(s["visual"]) or slide_figure_ref(s["title"])
                if fig_num is not None and caption_pages:
                    img_path = resolve_image_for_figure(
                        fig_num, caption_pages, image_files,
                        placed_via_caption, figure_to_image,
                    )
                    if img_path is not None:
                        cap_used += 1
                        image_source_note = (
                            f"[image: Figure {fig_num} caption page {caption_pages.get(fig_num)} "
                            f"→ {img_path.name}]"
                        )

                # Sequential fallback
                if img_path is None:
                    img_path = next_sequential_image()
                    if img_path is not None:
                        image_source_note = (
                            f"[image: sequential fallback → {img_path.name}]"
                        )

                if img_path is not None:
                    embed_path = maybe_resize(img_path, max_width, tmp_dir)
                    w, h = fit_image_box(embed_path, IMG_MAX_W, IMG_MAX_H)
                    cx = IMG_X + (IMG_MAX_W - w) / 2
                    cy = IMG_Y + (IMG_MAX_H - h) / 2
                    slide.shapes.add_picture(
                        str(embed_path),
                        Inches(cx), Inches(cy),
                        width=Inches(w), height=Inches(h),
                    )
                    images_used += 1
                else:
                    image_source_note = "[NEEDS IMAGE — out of extracted figures]"

        # ---- Speaker notes ----
        notes_parts = []
        if s["visual"]:
            notes_parts.append(f"[VISUAL] {s['visual']}")
        if image_source_note:
            notes_parts.append(image_source_note)
        if s["notes"]:
            notes_parts.append(s["notes"])
        slide.notes_slide.notes_text_frame.text = "\n\n".join(notes_parts)

    prs.save(str(out_path))
    size_mb = out_path.stat().st_size / (1024 * 1024)
    print(f"Wrote: {out_path}  ({size_mb:.1f} MB)")
    print(f"Slides: {len(slides)}")
    print(f"Images embedded: {images_used} / {images_total} available "
          f"(caption-paired: {cap_used}, sequential: {seq_used})")
    if max_width > 0:
        print(f"Image downsampling: max width = {max_width}px")
    return {
        "slides": len(slides),
        "images_total": images_total,
        "images_used": images_used,
        "caption_used": cap_used,
        "sequential_used": seq_used,
        "size_mb": size_mb,
    }


# ---------- CLI ---------------------------------------------------------------

def main():
    ap = argparse.ArgumentParser(
        description="Build PPTX from Stage 5 slide markdown with embedded figures.")
    ap.add_argument("paper_dir", help="Path to journal-club paper directory")
    ap.add_argument("--output", help="Output .pptx path (default: <paper_dir>/<paper_id>.pptx)")
    ap.add_argument("--images", help="Override images directory")
    ap.add_argument("--md", help="Override slide markdown path")
    ap.add_argument("--pdf", help="Override PDF path (for caption-aware figure mapping)")
    ap.add_argument("--max-width", type=int, default=1600,
                    help="Max image width in pixels; resizes larger images. "
                         "0 disables resizing. Default 1600.")
    ap.add_argument("--no-resize", action="store_true",
                    help="Disable image downsampling (equivalent to --max-width 0).")
    args = ap.parse_args()

    paper_dir = Path(args.paper_dir).expanduser().resolve()
    if not paper_dir.is_dir():
        sys.exit(f"Not a directory: {paper_dir}")

    md_path = Path(args.md) if args.md else paper_dir / "05_slides_draft.md"
    if not md_path.exists():
        sys.exit(f"Slide draft not found: {md_path}")

    # paper_id
    meta_path = paper_dir / "_meta.json"
    if meta_path.exists():
        meta = json.loads(meta_path.read_text())
        paper_id = meta.get("paper_id", paper_dir.name)
    else:
        paper_id = paper_dir.name

    # images dir: paper_dir/images, then paper_dir/<paper_id>/images
    if args.images:
        images_dir = Path(args.images).expanduser().resolve()
    else:
        candidates = [paper_dir / "images", paper_dir / paper_id / "images"]
        images_dir = next((p for p in candidates if p.exists()), candidates[0])

    # PDF discovery (for caption-aware figure mapping)
    if args.pdf:
        pdf_path = Path(args.pdf).expanduser().resolve()
    else:
        pdf_candidates: list[Path] = [paper_dir / "pdf" / f"{paper_id}.pdf"]
        if (paper_dir / "pdf").exists():
            pdf_candidates.extend(sorted((paper_dir / "pdf").glob("*.pdf")))
        pdf_candidates.extend(sorted(paper_dir.glob("*.pdf")))
        pdf_path = next((p for p in pdf_candidates if p.exists()), None)

    out_path = Path(args.output) if args.output else paper_dir / f"{paper_id}.pptx"
    max_width = 0 if args.no_resize else args.max_width

    print(f"Source markdown: {md_path}")
    print(f"Images dir:      {images_dir}{' (missing)' if not images_dir.exists() else ''}")
    print(f"PDF (captions):  {pdf_path or '(not found — using sequential image fallback)'}")
    print(f"Output:          {out_path}")
    print()

    slides = parse_slides(md_path)
    print(f"Parsed {len(slides)} slides.")
    build_pptx(slides, images_dir, out_path, pdf_path=pdf_path, max_width=max_width)


if __name__ == "__main__":
    main()
